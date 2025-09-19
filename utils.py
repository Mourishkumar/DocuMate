#import re
#import os
#import csv
#import json
#import requests
#import config as cfg
#from docx import Document
#from PyPDF2 import PdfReader
#from bs4 import BeautifulSoup
#from pptx import Presentation
#from dotenv import load_dotenv
#from typing import Optional, Dict
#from openpyxl import load_workbook
#
#from sentence_transformers import SentenceTransformer
#import easyocr
#import cv2
#
#ocr_reader = easyocr.Reader(['en'], gpu=cfg.DEVICE != 'cpu')
#MIN_REGION_AREA = 800
#DILATION_KERNEL_SIZE = (25, 5)
#
#def detect_text_regions(image) -> list[tuple[int, int, int, int]]:
#    """Detect text regions in an image using contour-based segmentation."""
#    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
#
#    _, binary_image = cv2.threshold(
#        gray_image, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU
#    )
#
#    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, DILATION_KERNEL_SIZE)
#    dilated_image = cv2.dilate(binary_image, kernel, iterations=1)
#
#    contours, _ = cv2.findContours(
#        dilated_image, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
#    )
#
#    bounding_boxes = []
#    img_height, img_width = binary_image.shape[:2]
#
#    for contour in contours:
#        x, y, w, h = cv2.boundingRect(contour)
#        region_area = w * h
#        if region_area < MIN_REGION_AREA:
#            continue
#
#        x0, y0 = max(0, x), max(0, y)
#        x1, y1 = min(img_width - 1, x + w), min(img_height - 1, y + h)
#        # bounding_boxes.append(image[y0:y1, x0:x1])
#        bounding_boxes.append((x0, y0, x1, y1))
#    bounding_boxes.sort(key=lambda box: (box[1], box[0]))
#    images = []
#    for (x0, y0, x1, y1) in bounding_boxes:
#        images.append(image[y0:y1, x0:x1])
#    return images
#
#def ocr(filepath):
#    image = cv2.imread(str(filepath))
#    data = []
#    images = detect_text_regions(image)
#    for segment in images:
#        for _, text, _ in ocr_reader.readtext(segment, detail=1):
#            data.append({
#                "text": text,
#            })
#    return data
#load_dotenv()
#
#if cfg.user_local_embeddings:
#    print('╭───────────────────────────────╮')
#    print('│ Using embedding model locally │')
#    print('╰───────────────────────────────╯')
#    _model = SentenceTransformer("BAAI/bge-m3", device=cfg.DEVICE, trust_remote_code=True)
#else:
#    _model = None
#
#def extract_thought(text: str) -> list[str]:
#    pattern = r'<think>(.*?)</think>'
#    return re.findall(pattern, text, re.DOTALL)
#
#def extract_text_from_file(filepath):
#    ext = os.path.splitext(filepath)[1].lower()
#    pages = []
#
#    if ext == ".pdf":
#        with open(filepath, "rb") as file:
#            reader = PdfReader(file)
#            for i, page in enumerate(reader.pages):
#                text = page.extract_text()
#                if text:
#                    pages.append(text.strip())
#    
#    elif ext == ".txt":
#        with open(filepath, "r", encoding="utf-8") as f:
#            text = f.read()
#            pages.append(text.strip())
#    
#    elif ext == '.docx':
#        doc = Document(filepath)
#        for i, para in enumerate(doc.paragraphs):
#            text = para.text.strip()
#            if text:
#                pages.append(text)
#    
#    elif ext in [".html", ".htm"]:
#        with open(filepath, "r", encoding="utf-8") as f:
#            soup = BeautifulSoup(f, "html.parser")
#            text = soup.get_text(separator="\n")
#            if text:
#                pages.append(text.strip())
#    
#    elif ext in [".csv", ".tsv"]:
#        delimiter = '\t' if ext == ".tsv" else ','
#        with open(filepath, "r", encoding="utf-8") as f:
#            reader = csv.reader(f, delimiter=delimiter)
#            for row in reader:
#                if any(cell.strip() for cell in row):
#                    pages.append(", ".join(row))
#    
#    elif ext == ".json":
#        with open(filepath, "r", encoding="utf-8") as f:
#            data = json.load(f)
#            text = json.dumps(data, indent=2, ensure_ascii=False)
#            pages.append(text)
#    
#    elif ext == ".xlsx":
#        workbook = load_workbook(filename=filepath, data_only=True)
#        sheets = []
#        for sheet in workbook.worksheets:
#            temp = []
#            for row in sheet.iter_rows(values_only=True):
#                row_text = [str(cell).strip() for cell in row if cell is not None]
#                if row_text:
#                    temp.append("| ".join(row_text))
#            pages.append('\n'.join(temp))
#        
#    
#    elif ext == ".pptx":
#        presentation = Presentation(filepath)
#        for slide in presentation.slides:
#            slide_text = []
#            for shape in slide.shapes:
#                if hasattr(shape, "text"):
#                    text = shape.text.strip()
#                    if text:
#                        slide_text.append(text)
#            if slide_text:
#                pages.append("\n".join(slide_text))
#
#    elif ext in ['.png', '.jpg']:
#        return ['\n'.join([t['text'] for t in ocr(filepath)])]
#
#    else:
#        raise ValueError(f"Unsupported file type: {ext}")
#    return pages
#
#def remove_comments(string):
#    string = re.sub(r'//.*', '', string)
#    string = re.sub(r'#.*', '', string)
#    string = re.sub(r'/\*.*?\*/', '', string, flags=re.DOTALL)
#    return string.strip()
#
#def fix_trailing_commas(json_string: str) -> str:
#    json_string = re.sub(r',\s*}', '}', json_string)
#    json_string = re.sub(r',\s*]', ']', json_string)
#    return json_string
#
#def remove_chinese_characters(text: str) -> str:
#    pattern = re.compile(r'[\u4e00-\u9fff]')
#    return pattern.sub('', text)
#
#def extract_json(text: str) -> Optional[Dict]:
#    text = remove_comments(text)
#    if '```json' not in text:
#        text = f'```json\n{text}'
#    """Extract JSON from LLM output between triple backticks."""
#    match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL)
#    if match:
#        try:
#            json_string = match.group(1)
#            json_string = fix_trailing_commas(json_string)
#            json_string = json.loads(json_string)
#            return json_string
#        except json.JSONDecodeError as e:
#            return None
#    return None
#
#def extract_plaintext(text: str) -> Optional[Dict]:
#    """Extract plaintext from LLM output between triple backticks."""
#    match = re.search(r'```plaintext\s*(.*?)\s*```', text, re.DOTALL)
#    if match:
#        try:
#            plaintext_string = match.group(1)
#            return plaintext_string
#        except:
#            return text.replace('```plaintext', '')
#    return text.replace('```plaintext', '')
#
#def ollama_llm(uri, json_payload, auth_type=None, auth_key=None, stream=True):  
#    headers = None
#    if auth_type == "bearer_token":
#        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {auth_key}"}
#    elif auth_type == "api_key":
#        if auth_key is None:
#            raise ValueError(f"auth_type is {auth_type}, but auth_key is null")
#        headers = {"Content-Type": "application/json", "api-key": f"{auth_key}"}
#    json_payload['stream'] = stream
#    with requests.post(uri, json=json_payload, headers=headers, stream=stream) as response:
#        if response.status_code == 200:
#            for chunk in response.iter_lines(decode_unicode=True):
#                data = json.loads(chunk)
#                yield json.dumps({"response": data['response']}) + "\n" 
#        else:
#            yield json.dumps({"error": response.status_code, "message": response.text}) + "\n"
#
#
#def ollama_embed(uri, json_payload, auth_type, auth_key):  
#    headers = None
#    if auth_type == "bearer_token":
#        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {auth_key}"}
#    elif auth_type == "api_key":
#        if auth_key is None:
#            raise ValueError(f"auth_type is {auth_type}, but auth_key is null")
#        headers = {"Content-Type": "application/json", "api-key": f"{auth_key}"}
#    response = requests.post(uri, json=json_payload, headers=headers)
#
#    return response.json()
#
#def get_embeds(sentence):
#    if _model is None:
#        return ollama_embed(f"{cfg.OLLAMA_HOST}/api/embed", {'model': 'bge-m3', 'input': sentence}, None, None)['embeddings']
#    else:
#        if isinstance(sentence, str):
#            sentence = [sentence]
#        embeddings = _model.encode(sentence, convert_to_numpy=True)  # shape: (n, d)
#        return embeddings.tolist()
###############################################

import re
import os
import csv
import json
import requests
import config as cfg
from docx import Document
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
from pptx import Presentation
from dotenv import load_dotenv
from typing import Optional, Dict
from openpyxl import load_workbook

from sentence_transformers import SentenceTransformer
import easyocr
import cv2

ocr_reader = easyocr.Reader(['en'], gpu=cfg.DEVICE != 'cpu')
MIN_REGION_AREA = 800
DILATION_KERNEL_SIZE = (25, 5)

def detect_text_regions(image) -> list[tuple[int, int, int, int]]:
    """Detect text regions in an image using contour-based segmentation."""
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

    _, binary_image = cv2.threshold(
        gray_image, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU
    )

    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, DILATION_KERNEL_SIZE)
    dilated_image = cv2.dilate(binary_image, kernel, iterations=1)

    contours, _ = cv2.findContours(
        dilated_image, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE
    )

    bounding_boxes = []
    img_height, img_width = binary_image.shape[:2]

    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        region_area = w * h
        if region_area < MIN_REGION_AREA:
            continue

        x0, y0 = max(0, x), max(0, y)
        x1, y1 = min(img_width - 1, x + w), min(img_height - 1, y + h)
        bounding_boxes.append((x0, y0, x1, y1))
    bounding_boxes.sort(key=lambda box: (box[1], box[0]))
    images = []
    for (x0, y0, x1, y1) in bounding_boxes:
        images.append(image[y0:y1, x0:x1])
    return images

def ocr(filepath):
    image = cv2.imread(str(filepath))
    data = []
    images = detect_text_regions(image)
    for segment in images:
        for _, text, _ in ocr_reader.readtext(segment, detail=1):
            data.append({
                "text": text,
            })
    return data

load_dotenv()

if cfg.user_local_embeddings:
    print('╭───────────────────────────────╮')
    print('│ Using embedding model locally │')
    print('╰───────────────────────────────╯')
    _model = SentenceTransformer("BAAI/bge-m3", device=cfg.DEVICE, trust_remote_code=True)
else:
    _model = None

def extract_thought(text: str) -> list[str]:
    pattern = r'<think>(.*?)</think>'
    return re.findall(pattern, text, re.DOTALL)

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    pages = []

    if ext == ".pdf":
        with open(filepath, "rb") as file:
            reader = PdfReader(file)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    pages.append(text.strip())
    
    elif ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            text = f.read()
            pages.append(text.strip())
    
    elif ext == '.docx':
        doc = Document(filepath)
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                pages.append(text)
    
    elif ext in [".html", ".htm"]:
        with open(filepath, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            text = soup.get_text(separator="\n")
            if text:
                pages.append(text.strip())
    
    elif ext in [".csv", ".tsv"]:
        delimiter = '\t' if ext == ".tsv" else ','
        with open(filepath, "r", encoding="utf-8") as f:
            reader = csv.reader(f, delimiter=delimiter)
            for row in reader:
                if any(cell.strip() for cell in row):
                    pages.append(", ".join(row))
    
    elif ext == ".json":
        with open(filepath, "r", encoding="utf-8") as f:
            data = json.load(f)
            text = json.dumps(data, indent=2, ensure_ascii=False)
            pages.append(text)
    
    elif ext == ".xlsx":
        workbook = load_workbook(filename=filepath, data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            temp = []
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell).strip() for cell in row if cell is not None]
                if row_text:
                    temp.append("| ".join(row_text))
            pages.append('\n'.join(temp))
        
    elif ext == ".pptx":
        presentation = Presentation(filepath)
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                pages.append("\n".join(slide_text))

    elif ext in ['.png', '.jpg']:
        return ['\n'.join([t['text'] for t in ocr(filepath)])]

    else:
        raise ValueError(f"Unsupported file type: {ext}")
    return pages

def remove_comments(string):
    string = re.sub(r'//.*', '', string)
    string = re.sub(r'#.*', '', string)
    string = re.sub(r'/\*.*?\*/', '', string, flags=re.DOTALL)
    return string.strip()

def fix_trailing_commas(json_string: str) -> str:
    json_string = re.sub(r',\s*}', '}', json_string)
    json_string = re.sub(r',\s*]', ']', json_string)
    return json_string

def remove_chinese_characters(text: str) -> str:
    pattern = re.compile(r'[\u4e00-\u9fff]')
    return pattern.sub('', text)

def extract_json(text: str) -> Optional[Dict]:
    text = remove_comments(text)
    if '```json' not in text:
        text = f'```json\n{text}'
    """Extract JSON from LLM output between triple backticks."""
    match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL)
    if match:
        try:
            json_string = match.group(1)
            json_string = fix_trailing_commas(json_string)
            json_string = json.loads(json_string)
            return json_string
        except json.JSONDecodeError as e:
            return None
    return None

def extract_plaintext(text: str) -> Optional[Dict]:
    """Extract plaintext from LLM output between triple backticks."""
    match = re.search(r'```plaintext\s*(.*?)\s*```', text, re.DOTALL)
    if match:
        try:
            plaintext_string = match.group(1)
            return plaintext_string
        except:
            return text.replace('```plaintext', '')
    return text.replace('```plaintext', '')

def ollama_llm(uri, json_payload, auth_type=None, auth_key=None, stream=True):  
    headers = None
    if auth_type == "bearer_token":
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {auth_key}"}
    elif auth_type == "api_key":
        if auth_key is None:
            raise ValueError(f"auth_type is {auth_type}, but auth_key is null")
        headers = {"Content-Type": "application/json", "api-key": f"{auth_key}"}
    json_payload['stream'] = stream
    with requests.post(uri, json=json_payload, headers=headers, stream=stream) as response:
        if response.status_code == 200:
            for chunk in response.iter_lines(decode_unicode=True):
                data = json.loads(chunk)
                yield json.dumps({"response": data['response']}) + "\n" 
        else:
            yield json.dumps({"error": response.status_code, "message": response.text}) + "\n"


def ollama_embed(uri, json_payload, auth_type, auth_key):  
    headers = None
    if auth_type == "bearer_token":
        headers = {"Content-Type": "application/json", "Authorization": f"Bearer {auth_key}"}
    elif auth_type == "api_key":
        if auth_key is None:
            raise ValueError(f"auth_type is {auth_type}, but auth_key is null")
        headers = {"Content-Type": "application/json", "api-key": f"{auth_key}"}
    response = requests.post(uri, json=json_payload, headers=headers)
    return response.json()

def get_embeds(sentence):
    if _model is None:  # use Ollama remote embeddings
        return ollama_embed(
            f"{cfg.OLLAMA_HOST}/api/embed",
            {'model': cfg.OLLAMA_EMBED_MODEL, 'input': sentence},
            None,
            None
        )['embeddings']
    else:  # use local SentenceTransformer
        if isinstance(sentence, str):
            sentence = [sentence]
        embeddings = _model.encode(sentence, convert_to_numpy=True)  # shape: (n, d)
        return embeddings.tolist()
